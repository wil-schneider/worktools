import collections 
import collections.abc
from pptx import Presentation
from PIL import Image
import io
import os

def extract_images_from_pptx(pptx_file, target_dir):
    # Load the Presentation from the pptx file
    presentation = Presentation(pptx_file)

    # Go through every slide
    for slide_no, slide in enumerate(presentation.slides):
        # Go through every shape in the slide
        for shape_no, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                # Get image data
                image = shape.image
                image_bytes = image.blob

                # Create PIL image
                pil_image = Image.open(io.BytesIO(image_bytes))

                # Create target directory if it doesn't exist
                if not os.path.exists(target_dir):
                    os.makedirs(target_dir)

                # Save the image
                image_filename = f"slide_{slide_no+1}_shape_{shape_no+1}.png"
                pil_image.save(os.path.join(target_dir, image_filename))

# Use function on a given pptx file
extract_images_from_pptx("MDM Intermediate Quizzes.pptx", "pptOutputDir")
